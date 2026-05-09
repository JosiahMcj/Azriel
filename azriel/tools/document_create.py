"""document_create tool -- generate .docx / .pptx / .xlsx files.

Saves the result inside the sandbox at ~/azriel-files/. Returns the
relative path so the user can fetch it later.

Format string:
  document_create("docx|sermon-2026-05-04|<TITLE>\\n\\n<paragraphs separated by blank lines>")
  document_create("xlsx|attendance|Date,Count\\n2026-05-04,42\\n2026-05-11,55") -- CSV body
  document_create("pptx|outline|<TITLE>\\n---\\nSlide 1 title\\nbullet 1\\nbullet 2\\n---\\nSlide 2 title\\nbullet")
"""
from datetime import datetime
from pathlib import Path

SANDBOX = Path.home() / "azriel-files"


def _safe_name(name: str) -> str:
    keep = "".join(c if c.isalnum() or c in "-_." else "_" for c in name)
    return keep[:80] or "untitled"


def _make_docx(title: str, body: str, out: Path):
    try:
        from docx import Document
    except ImportError:
        return "ERROR: python-docx not installed. uv pip install python-docx"
    doc = Document()
    if title:
        doc.add_heading(title, level=1)
    for para in body.split("\n\n"):
        para = para.strip()
        if para:
            doc.add_paragraph(para)
    doc.save(str(out))
    return None


def _make_xlsx(title: str, body: str, out: Path):
    try:
        from openpyxl import Workbook
    except ImportError:
        return "ERROR: openpyxl not installed. uv pip install openpyxl"
    wb = Workbook()
    ws = wb.active
    ws.title = (title or "Sheet1")[:31]
    for row in body.split("\n"):
        if not row.strip():
            continue
        ws.append([c.strip() for c in row.split(",")])
    wb.save(str(out))
    return None


def _make_pptx(title: str, body: str, out: Path):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return "ERROR: python-pptx not installed. uv pip install python-pptx"
    prs = Presentation()
    sections = body.split("\n---\n")
    # First section is the title slide
    title_layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = title or "Untitled"
    if len(sections[0].strip()) > 0 and sections[0].strip() != title:
        try:
            s.placeholders[1].text = sections[0].strip()
        except (KeyError, IndexError):
            pass
    for sec in sections[1:]:
        lines = [ln for ln in sec.strip().split("\n") if ln.strip()]
        if not lines:
            continue
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = lines[0]
        body_ph = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body_ph and len(lines) > 1:
            tf = body_ph.text_frame
            tf.text = lines[1]
            for line in lines[2:]:
                p = tf.add_paragraph()
                p.text = line
    prs.save(str(out))
    return None


def document_create(arg: str) -> str:
    if not isinstance(arg, str):
        return "ERROR: document_create expects 'format|name|content'."
    parts = arg.split("|", 2)
    if len(parts) < 3:
        return "ERROR: format is 'format|name|content' (format ∈ docx/pptx/xlsx)."
    fmt, name, content = parts[0].strip().lower(), parts[1].strip(), parts[2]
    if fmt not in ("docx", "pptx", "xlsx"):
        return f"ERROR: unsupported format '{fmt}'. Use docx, pptx, or xlsx."
    SANDBOX.mkdir(parents=True, exist_ok=True)
    safe = _safe_name(name) + "." + fmt
    out = SANDBOX / safe
    title = name.split("\n", 1)[0]
    err = (_make_docx if fmt == "docx" else _make_xlsx if fmt == "xlsx" else _make_pptx)(
        title, content, out
    )
    if err:
        return err
    size = out.stat().st_size
    return f"created {safe} ({size:,} bytes) -- accessible via fs_read or download from {out}"


if __name__ == "__main__":
    import sys
    print(document_create(sys.argv[1] if len(sys.argv) > 1 else "docx|test|Hello\n\nWorld"))
